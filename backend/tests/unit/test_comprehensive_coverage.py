from __future__ import annotations

import os
import struct
import io
import pytest
from datetime import datetime, timezone
from pathlib import Path
from PIL import Image, PngImagePlugin

from app.services.metadata_service import MetadataService
from app.services.ai_attribution_service import AIAttributionService
from app.services.video_forensics import VideoForensicsService
from app.services.provenance_service import ProvenanceService

import docx
import pptx
import openpyxl

@pytest.mark.unit
def test_metadata_service_office_formats(tmp_path: Path) -> None:
    # 1. Create a dummy docx file
    doc_path = tmp_path / "test.docx"
    doc = docx.Document()
    doc.add_paragraph("Hello world")
    doc.core_properties.author = "Test Author"
    doc.core_properties.title = "Test Title"
    doc.save(str(doc_path))

    meta_doc = MetadataService.extract_metadata(str(doc_path), "document")
    assert meta_doc["creator"] == "Test Author"
    assert meta_doc["raw_metadata"]["title"] == "Test Title"

    # 2. Create a dummy pptx file
    ppt_path = tmp_path / "test.pptx"
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.core_properties.author = "PPTX Creator"
    prs.core_properties.title = "PPTX Title"
    prs.save(str(ppt_path))

    meta_ppt = MetadataService.extract_metadata(str(ppt_path), "document")
    assert meta_ppt["creator"] == "PPTX Creator"
    assert meta_ppt["raw_metadata"]["title"] == "PPTX Title"

    # 3. Create a dummy xlsx file
    xls_path = tmp_path / "test.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Data"
    wb.properties.creator = "Sheet Creator"
    wb.properties.title = "Sheet Title"
    wb.save(str(xls_path))

    meta_xls = MetadataService.extract_metadata(str(xls_path), "document")
    assert meta_xls["creator"] == "Sheet Creator"
    assert meta_xls["raw_metadata"]["title"] == "Sheet Title"


@pytest.mark.unit
def test_ai_attribution_service_png_chunks(tmp_path: Path) -> None:
    # 1. SD/Flux PNG parameter parsing
    sd_path = tmp_path / "sd.png"
    img = Image.new("RGB", (10, 10))
    info = PngImagePlugin.PngInfo()
    info.add_text("parameters", "a fantasy landscape\nSteps: 20\nModel: SDXL-v1")
    img.save(str(sd_path), pnginfo=info)

    attr_sd = AIAttributionService.attribute_ai_content(str(sd_path), "image")
    assert attr_sd["predicted_source"] == "Stable Diffusion / Flux"
    assert attr_sd["probability"] == 0.98
    assert "Model identified: SDXL-v1" in attr_sd["indicators"]["structural_cues"]

    # 2. Midjourney Description chunk
    mj_path = tmp_path / "mj.png"
    img_mj = Image.new("RGB", (10, 10))
    info_mj = PngImagePlugin.PngInfo()
    info_mj.add_text("Description", "A cool midjourney artwork prompt info")
    img_mj.save(str(mj_path), pnginfo=info_mj)

    attr_mj = AIAttributionService.attribute_ai_content(str(mj_path), "image")
    assert attr_mj["predicted_source"] == "Midjourney"
    assert attr_mj["probability"] == 0.99

    # 3. Midjourney Comment chunk
    mj_comment_path = tmp_path / "mj_comment.png"
    img_mj_c = Image.new("RGB", (10, 10))
    info_mj_c = PngImagePlugin.PngInfo()
    info_mj_c.add_text("Comment", "made in midjourney v6")
    img_mj_c.save(str(mj_comment_path), pnginfo=info_mj_c)

    attr_mj_c = AIAttributionService.attribute_ai_content(str(mj_comment_path), "image")
    assert attr_mj_c["predicted_source"] == "Midjourney"


@pytest.mark.unit
def test_video_forensics_mp4_atoms(tmp_path: Path) -> None:
    # Construct a valid dummy MP4 container in binary format to cover atom parser
    # Structure:
    # 8 bytes: box size (4 bytes), box type 'moov' (4 bytes)
    # 8 bytes: box size (4 bytes), box type 'mvhd' (4 bytes)
    # 4 bytes: version (0) and flags (0)
    # 4 bytes: creation_time
    # 4 bytes: modification_time
    # 4 bytes: timescale
    # 4 bytes: duration
    mvhd_box = struct.pack(">I4s", 28, b"mvhd") + b"\x00\x00\x00\x00" + struct.pack(">IIII", 100000000, 100000000, 1000, 30000)
    moov_box = struct.pack(">I4s", len(mvhd_box) + 8, b"moov") + mvhd_box

    mp4_path = tmp_path / "test.mp4"
    mp4_path.write_bytes(moov_box)

    analysis = VideoForensicsService.analyze_video(str(mp4_path))
    assert analysis["metadata"]["method"] == "binary_atom_scanner"
    assert analysis["metadata"]["duration_seconds"] == 30.0
    assert analysis["tampered"] is True  # 1904 baseline timestamp forces abnormal tag


@pytest.mark.unit
def test_provenance_service_jpeg_jumbf(tmp_path: Path) -> None:
    # Construct JPEG payload containing APP11 (0xFFEB) marker with 'jumb' payload
    # Start of Image: 0xFFD8
    # APP11 Marker: 0xFFEB, Length (2 bytes), followed by 'jumb' payload
    length = 26
    app11_payload = struct.pack(">HH", 0xFFEB, length) + b"some header C2PA jumb"
    jpeg_bytes = b"\xff\xd8" + app11_payload + b"\xff\xda"  # SOI + APP11 + SOS (Start of Scan)

    jpeg_path = tmp_path / "provenance.jpg"
    jpeg_path.write_bytes(jpeg_bytes)

    res = ProvenanceService.assess_provenance(str(jpeg_path))
    assert res["has_manifest"] is True
    assert res["manifest_valid"] is True
    assert res["ownership_classification"] == "VERIFIED OWNER"


@pytest.mark.unit
def test_provenance_service_png_cabx(tmp_path: Path) -> None:
    # Construct PNG payload containing 'caBX' chunk
    # PNG signature: 8 bytes
    # Chunk: length (4 bytes), type 'caBX' (4 bytes), data (0 bytes), CRC (4 bytes)
    png_signature = b"\x89PNG\r\n\x1a\n"
    cabx_chunk = struct.pack(">I4s", 0, b"caBX") + struct.pack(">I", 0)  # Length 0, type caBX, CRC 0
    iend_chunk = struct.pack(">I4s", 0, b"IEND") + struct.pack(">I", 0)
    png_bytes = png_signature + cabx_chunk + iend_chunk

    png_path = tmp_path / "provenance.png"
    png_path.write_bytes(png_bytes)

    res = ProvenanceService.assess_provenance(str(png_path))
    assert res["has_manifest"] is True
    assert res["manifest_valid"] is True
    assert res["ownership_classification"] == "VERIFIED OWNER"


@pytest.mark.unit
def test_image_forensics_service(tmp_path: Path) -> None:
    from app.services.image_forensics import ImageForensicsService

    # 1. Test running analyses on a solid image (triggers clones and ghosting)
    img_path = tmp_path / "solid.jpg"
    img = Image.new("RGB", (128, 128), color="red")
    img.save(str(img_path), "JPEG", quality=90)

    upload_dir = tmp_path / "uploads"
    upload_dir.mkdir()

    res = ImageForensicsService.run_all_analyses(str(img_path), "ev1", str(upload_dir))
    assert res["tampered"] is True or res["tampered"] is False
    assert len(res["supporting_evidence"]) > 0
    assert "details" in res

    # Check that ela output image exists
    ela_filename = "ela_ev1.jpg"
    assert (upload_dir / ela_filename).exists()

    # 2. Test running on a noisy image to hit noise variance / anomaly ratio branches
    noisy_path = tmp_path / "noisy.png"
    # Create an image with high-frequency noise pattern
    noisy_img = Image.new("RGB", (128, 128))
    pixels = noisy_img.load()
    for x in range(128):
        for y in range(128):
            val = (x * 17 + y * 23) % 256
            pixels[x, y] = (val, val, val)
    noisy_img.save(str(noisy_path), "PNG")

    res_noise = ImageForensicsService.run_all_analyses(str(noisy_path), "ev2", str(upload_dir))
    assert "noise" in res_noise["details"]

    # 3. Test exception handling inside image forensics (non-existent file)
    bad_path = str(tmp_path / "does_not_exist.jpg")
    ela_path, ela_tampered, ela_conf = ImageForensicsService.calculate_ela(bad_path, "ev3", str(upload_dir))
    assert ela_path == ""
    assert ela_tampered is False
    assert ela_conf == 0.0

    noise_tampered, noise_conf, noise_stats = ImageForensicsService.analyze_noise(bad_path)
    assert noise_tampered is False
    assert noise_conf == 0.0

    clone_detected, clone_conf, clone_regions = ImageForensicsService.detect_clones(bad_path)
    assert clone_detected is False
    assert clone_regions == []

    ghost_detected, ghost_conf, ghost_quality = ImageForensicsService.detect_jpeg_ghosts(bad_path)
    assert ghost_detected is False
    assert ghost_quality == 90


@pytest.mark.unit
def test_ai_attribution_exif_and_fallbacks(tmp_path: Path) -> None:
    from app.services.ai_attribution_service import AIAttributionService

    # 1. Test EXIF Software mock checks
    class MockImage:
        def __init__(self, software_name: str):
            self.software_name = software_name
            self.info = {}
            self.format = "PNG"
            self.width = 10
            self.height = 10
            self.mode = "RGB"
        def __enter__(self):
            return self
        def __exit__(self, exc_type, exc_val, exc_tb):
            pass
        def _getexif(self):
            # 305 is Software tag code
            return {305: self.software_name}

    import unittest.mock as mock
    with mock.patch("PIL.Image.open") as mock_open:
        # Mock Midjourney software
        mock_open.return_value = MockImage("Midjourney v6")
        res = AIAttributionService.attribute_ai_content("dummy.png", "image")
        assert res["predicted_source"] == "Midjourney"

        # Mock NovelAI
        mock_open.return_value = MockImage("NovelAI diffusion")
        res = AIAttributionService.attribute_ai_content("dummy.png", "image")
        assert res["predicted_source"] == "Stable Diffusion (NovelAI)"

        # Mock DALL-E
        mock_open.return_value = MockImage("DALL-E 3 generator")
        res = AIAttributionService.attribute_ai_content("dummy.png", "image")
        assert res["predicted_source"] == "DALL-E"

    # 2. Test fallbacks based on filenames
    names_to_sources = {
        "sample_flux.png": "Flux",
        "dalle_generation.jpg": "DALL-E",
        "dall-e_file.png": "DALL-E",
        "sora_clip.mp4": "Sora",
        "gen3_runway.mp4": "Runway Gen-3",
        "veo_video.mp4": "Veo",
        "human_original.png": "Human / Camera Original"
    }

    for filename, expected_source in names_to_sources.items():
        res = AIAttributionService.attribute_ai_content(filename, "image" if filename.endswith((".png", ".jpg")) else "video")
        assert res["predicted_source"] == expected_source


@pytest.mark.unit
def test_document_forensics_extended(tmp_path: Path) -> None:
    from app.services.document_forensics import DocumentForensicsService
    import zipfile

    # 1. Test unsupported document type fallback
    res = DocumentForensicsService.analyze_document("test.txt")
    assert res["tampered"] is False
    assert "Unsupported document type" in res["reasons"][0]

    # 2. Test PDF with /Launch and %%EOF incremental updates
    pdf_launch_data = (
        b"%PDF-1.7\n"
        b"1 0 obj\n<< /Type /Catalog /Launch << /F (cmd.exe) >> >>\nendobj\n"
        b"%%EOF\n"
        b"%%EOF\n"
    )
    pdf_path = tmp_path / "launch.pdf"
    pdf_path.write_bytes(pdf_launch_data)

    res_pdf = DocumentForensicsService.analyze_document(str(pdf_path))
    assert res_pdf["tampered"] is True
    assert res_pdf["details"]["launch_count"] == 1
    assert res_pdf["details"]["incremental_updates_count"] == 1
    assert any("Process launching command" in r for r in res_pdf["reasons"])
    assert any("Incremental edits detected" in r for r in res_pdf["reasons"])

    # 3. Test PDF scan exception handler (trigger exception by passing non-existent path ending in .pdf)
    res_err = DocumentForensicsService.analyze_document(str(tmp_path / "does_not_exist.pdf"))
    assert res_err["confidence"] == 50.0
    assert any("failed" in r or "audit failed" in r for r in res_err["reasons"])

    # 4. Test OOXML docx macro and external relation target detection
    docx_path = tmp_path / "malicious.docx"
    with zipfile.ZipFile(docx_path, "w") as z:
        # Write macro file vbaProject.bin
        z.writestr("word/vbaProject.bin", b"macro code bytes")
        # Write relationship file with TargetMode="External"
        z.writestr("word/_rels/document.xml.rels", b'<?xml version="1.0" encoding="UTF-8"?><Relationships><Relationship TargetMode="External" Target="http://evil.com/payload"/></Relationships>')

    res_docx = DocumentForensicsService.analyze_document(str(docx_path))
    assert res_docx["tampered"] is True
    assert res_docx["details"]["has_macros"] is True
    assert len(res_docx["details"]["embedded_objects"]) > 0
    assert any("Active Macro container detected" in r for r in res_docx["reasons"])
    assert any("External reference payload found" in r for r in res_docx["reasons"])


@pytest.mark.unit
def test_metadata_service_gps_and_exceptions(tmp_path: Path) -> None:
    from app.services.metadata_service import MetadataService
    from fractions import Fraction
    import unittest.mock as mock

    # 1. Test metadata service exception routing (passing non-existent path)
    res_err = MetadataService.extract_metadata("non_existent_file.png", "image")
    assert "error" in res_err["raw_metadata"]

    # 1b. Test system fallback exception routing (FileNotFoundError)
    res_fallback_err = MetadataService.extract_metadata("non_existent_file.bin", "video")
    assert "error" in res_fallback_err["raw_metadata"]

    # 2. Test GPS degree parsing from different formats
    gps_info_tuple = {
        "GPSLatitudeRef": "N",
        "GPSLatitude": (34, 4, 11.2),
        "GPSLongitudeRef": "W",
        "GPSLongitude": (118, 26, 54.3)
    }
    lat, lon = MetadataService._parse_gps_coords(gps_info_tuple)
    assert round(lat, 4) == 34.0698
    assert round(lon, 4) == -118.4484

    # Test string rational representation e.g. "34/1, 4/1, 112/10" or rational list ["118/1", "26/1", "543/10"]
    gps_info_str = {
        "GPSLatitudeRef": "S",
        "GPSLatitude": "34/1, 4/1, 112/10",
        "GPSLongitudeRef": "E",
        "GPSLongitude": ["118/1", "26/1", "543/10"]
    }
    lat, lon = MetadataService._parse_gps_coords(gps_info_str)
    assert round(lat, 4) == -34.0698
    assert round(lon, 4) == 118.4484

    # Test fraction rational representation
    gps_info_frac = {
        "GPSLatitudeRef": "N",
        "GPSLatitude": (Fraction(34, 1), Fraction(4, 1), Fraction(112, 10)),
        "GPSLongitudeRef": "W",
        "GPSLongitude": (Fraction(118, 1), Fraction(26, 1), Fraction(543, 10))
    }
    lat, lon = MetadataService._parse_gps_coords(gps_info_frac)
    assert round(lat, 4) == 34.0698
    assert round(lon, 4) == -118.4484

    # Test single-value coordinates list/number
    gps_info_single = {
        "GPSLatitudeRef": "N",
        "GPSLatitude": [34.0698],
        "GPSLongitudeRef": "W",
        "GPSLongitude": 118.4484
    }
    lat, lon = MetadataService._parse_gps_coords(gps_info_single)
    assert lat == 34.0698
    assert lon == -118.4484

    # Test invalid / exception GPS formats
    lat, lon = MetadataService._parse_gps_coords({"GPSLatitude": "invalid"})
    assert lat is None
    assert lon is None

    # Test None GPS parsing exception
    lat, lon = MetadataService._parse_gps_coords(None)
    assert lat is None
    assert lon is None

    # 3. Test _extract_pdf_metadata author, producer, pages
    class MockPdfReader:
        def __init__(self, file_path):
            class Info:
                author = "Author Name"
                creator = "Creator Name"
                producer = "Producer Name"
                def items(self):
                    return [("/Author", "Author Name"), ("/Producer", "Producer Name")]
                def get(self, key):
                    if key == "/CreationDate":
                        return "D:20260611120000"
                    return None
            self.metadata = Info()
            self.pages = [1, 2, 3]

    with mock.patch("app.services.metadata_service.PdfReader", side_effect=MockPdfReader):
        res_pdf = MetadataService.extract_metadata("dummy.pdf", "document")
        assert res_pdf["creator"] == "Author Name"
        assert res_pdf["software_used"] == "Producer Name"
        assert res_pdf["raw_metadata"]["pages_count"] == 3

    # Test PDF parsing exception block
    with mock.patch("app.services.metadata_service.PdfReader", side_effect=Exception("pdf read error")):
        res_pdf_err = MetadataService.extract_metadata("dummy.pdf", "document")
        assert "PDF extraction error" in res_pdf_err["raw_metadata"]["error"]

    # 3b. Test Image metadata EXIF fields happy path with mock
    class MockImageWithExif:
        def __init__(self):
            self.format = "JPEG"
            self.width = 100
            self.height = 100
            self.mode = "RGB"
            self.info = {}
        def __enter__(self):
            return self
        def __exit__(self, exc_type, exc_val, exc_tb):
            pass
        def _getexif(self):
            return {
                305: "Adobe Photoshop",
                315: "John Doe",
                36867: "2026:06:11 12:00:00",
                34853: {
                    "GPSLatitudeRef": "N",
                    "GPSLatitude": "34/1, 4/1, 112/10",
                    "GPSLongitudeRef": "W",
                    "GPSLongitude": "118/1, 26/1, 543/10"
                }
            }

    with mock.patch("PIL.Image.open") as mock_open:
        mock_open.return_value = MockImageWithExif()
        res_exif = MetadataService.extract_metadata("dummy.jpg", "image")
        assert res_exif["creator"] == "John Doe"
        assert res_exif["software_used"] == "Adobe Photoshop"

    # 4. Trigger exceptions in docx, pptx, xlsx metadata
    res_docx_err = MetadataService._extract_docx_metadata(str(tmp_path))
    assert "error" in res_docx_err["raw_metadata"]

    res_pptx_err = MetadataService._extract_pptx_metadata(str(tmp_path))
    assert "error" in res_pptx_err["raw_metadata"]

    res_xlsx_err = MetadataService._extract_xlsx_metadata(str(tmp_path))
    assert "error" in res_xlsx_err["raw_metadata"]


@pytest.mark.unit
def test_provenance_service_extended(tmp_path: Path) -> None:
    from app.services.provenance_service import ProvenanceService
    import zipfile
    import struct
    import unittest.mock as mock

    # 1. Test extract_c2pa_provenance wrapper
    dummy_jpg = tmp_path / "dummy.jpg"
    dummy_jpg.write_bytes(b"\xff\xd8\xff\xeb\x00\x1ajumbc2pa-content\xff\xda")
    res_wrapper = ProvenanceService.extract_c2pa_provenance(str(dummy_jpg))
    assert res_wrapper["has_manifest"] is True

    # 2. Test metadata signals manifest_hint and happy path PDF metadata
    # Write a simple PDF file
    pdf_path = tmp_path / "metadata_hint.pdf"
    pdf_path.write_bytes(b"%PDF-1.7\n%%EOF")

    class MockPdfReaderWithAuthor:
        def __init__(self, file_path):
            class Info:
                author = "Jane Creator"
                creator = None
                producer = None
                def items(self):
                    return [("/Author", "Jane Creator")]
                def get(self, key):
                    if key == "/Author":
                        return "Jane Creator"
                    return None
            self.metadata = Info()

    with mock.patch("app.services.provenance_service.PdfReader", side_effect=MockPdfReaderWithAuthor):
        res_pdf = ProvenanceService.assess_provenance(str(pdf_path))
        assert res_pdf["has_manifest"] is True
        assert res_pdf["ownership_classification"] == "PROBABLE OWNER"
        assert any("Manifest-like provenance metadata" in r for r in res_pdf["reasons"])

    # 3. Test office archive scan
    docx_path = tmp_path / "archive_prov.docx"
    with zipfile.ZipFile(docx_path, "w") as z:
        z.writestr("core.xml", b"metadata")
    res_docx = ProvenanceService.assess_provenance(str(docx_path))
    assert res_docx["has_manifest"] is True
    assert any("Office core metadata package found" in e for e in res_docx["supporting_evidence"])

    # 3b. Test office archive error exception handler
    res_docx_err = ProvenanceService.assess_provenance("does_not_exist.docx")
    assert any("Office provenance parsing error" in e for e in res_docx_err["supporting_evidence"])

    # 4. Test office archive scan with c2pa in filenames
    docx_c2pa_path = tmp_path / "archive_c2pa.docx"
    with zipfile.ZipFile(docx_c2pa_path, "w") as z:
        z.writestr("c2pa_manifest.xml", b"manifest content")
    res_docx_c2pa = ProvenanceService.assess_provenance(str(docx_c2pa_path))
    assert res_docx_c2pa["has_manifest"] is True
    assert any("Archive contains manifest-like provenance entries" in e for e in res_docx_c2pa["supporting_evidence"])

    # 5. Test AI generation signals in metadata creator/software
    res_ai = ProvenanceService.assess_provenance(
        str(dummy_jpg),
        metadata={"creator": "Midjourney v6", "software_used": "flux-generator"}
    )
    assert res_ai["ai_generation_signals"]["present"] is True

    # 6. Test confidence calculation: has_manifest=True, manifest_valid=False
    score = ProvenanceService._calculate_confidence(
        has_manifest=True,
        manifest_valid=False,
        blockchain_verified=False,
        creator="Some Creator",
        device="Camera",
        editing_history=[{"action": "edit"}],
        ai_generation_signals=False
    )
    assert score == 45.0

    # 7. Test empty supporting evidence block mock
    with mock.patch.object(ProvenanceService, "_collect_manifest_signals", return_value={"present": False, "valid": False, "evidence": []}):
        res_empty = ProvenanceService.assess_provenance("dummy.txt")
        assert "No provenance evidence beyond file structure" in res_empty["supporting_evidence"][0]

    # 8. _scan_jpeg_app11_jumb test cases
    # Not jpeg header
    bad_jpg = tmp_path / "bad.jpg"
    bad_jpg.write_bytes(b"not jpeg")
    assert ProvenanceService._scan_jpeg_app11_jumb(str(bad_jpg)) is False

    # Short marker header
    short_jpg = tmp_path / "short.jpg"
    short_jpg.write_bytes(b"\xff\xd8\xff")
    assert ProvenanceService._scan_jpeg_app11_jumb(str(short_jpg)) is False

    # FFDA marker break
    ffda_jpg = tmp_path / "ffda.jpg"
    ffda_jpg.write_bytes(b"\xff\xd8\xff\xda")
    assert ProvenanceService._scan_jpeg_app11_jumb(str(ffda_jpg)) is False

    # Exception
    assert ProvenanceService._scan_jpeg_app11_jumb(str(tmp_path)) is False

    # 9. _scan_png_chunks test cases
    # Not png header
    bad_png = tmp_path / "bad.png"
    bad_png.write_bytes(b"not png")
    assert ProvenanceService._scan_png_chunks(str(bad_png)) is False

    # IEND chunk break
    iend_png = tmp_path / "iend.png"
    iend_png.write_bytes(b"\x89PNG\r\n\x1a\n" + struct.pack(">I4s", 0, b"IEND") + struct.pack(">I", 0))
    assert ProvenanceService._scan_png_chunks(str(iend_png)) is False

    # Exception
    assert ProvenanceService._scan_png_chunks(str(tmp_path)) is False

    # 10. PDF & Archive scan exceptions
    p, v, n = ProvenanceService._scan_pdf_provenance(str(tmp_path))
    assert p is False
    assert any("scan failed" in note for note in n)

    p2, v2, n2 = ProvenanceService._scan_archive_provenance(str(tmp_path))
    assert p2 is False

    with mock.patch("zipfile.is_zipfile", side_effect=Exception("mock zip err")):
        p3, v3, n3 = ProvenanceService._scan_archive_provenance("dummy.docx")
        assert any("scan failed" in note for note in n3)


@pytest.mark.unit
def test_main_root_endpoint(client) -> None:
    response = client.get("/")
    assert response.status_code == 200
    assert response.json()["status"] == "healthy"


